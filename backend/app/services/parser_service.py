import hashlib
import os
import threading
import uuid
from pathlib import Path
from typing import List, Optional

from app.config import settings

_docling_converter_lock = threading.Lock()
_docling_converter = None

_easyocr_reader_lock = threading.Lock()
_easyocr_reader = None


def _get_docling_converter():
    global _docling_converter
    if _docling_converter is None:
        with _docling_converter_lock:
            if _docling_converter is None:
                try:
                    from docling.document_converter import DocumentConverter
                    _docling_converter = DocumentConverter()
                except ImportError:
                    _docling_converter = False
    return _docling_converter if _docling_converter is not False else None


def _get_easyocr_reader():
    global _easyocr_reader
    if _easyocr_reader is None:
        with _easyocr_reader_lock:
            if _easyocr_reader is None:
                import easyocr
                _easyocr_reader = easyocr.Reader(["ch_sim", "en"], gpu=True)
    return _easyocr_reader


class ParsedChunk:
    def __init__(self, content: str, page_number: Optional[int] = None, metadata: Optional[dict] = None):
        self.content = content
        self.page_number = page_number
        self.metadata = metadata or {}


class DoclingParser:

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".txt", ".html", ".htm", ".png", ".jpg", ".jpeg"}

    def __init__(self, ocr_enabled: bool = False):
        self.ocr_enabled = ocr_enabled

    def parse(self, file_path: str) -> List[ParsedChunk]:
        ext = Path(file_path).suffix.lower()

        if ext in {".png", ".jpg", ".jpeg"} and self.ocr_enabled:
            ocr_result = self._parse_image(file_path)
            if ocr_result and ocr_result[0].content and not ocr_result[0].content.startswith("[Image"):
                return ocr_result

        converter = _get_docling_converter()
        if converter:
            try:
                return self._parse_with_docling(file_path, converter)
            except Exception:
                pass

        return self._parse_basic(file_path, ext)

    def _parse_with_docling(self, file_path: str, converter=None) -> List[ParsedChunk]:
        if converter is None:
            converter = _get_docling_converter()
        result = converter.convert(file_path)
        text = result.document.export_to_markdown()
        return [ParsedChunk(content=text, metadata={"source": "docling"})]

    def _parse_basic(self, file_path: str, ext: str) -> List[ParsedChunk]:
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".docx":
            return self._parse_docx(file_path)
        elif ext == ".pptx":
            return self._parse_pptx(file_path)
        elif ext in {".md", ".txt"}:
            return self._parse_text(file_path)
        elif ext in {".html", ".htm"}:
            return self._parse_html(file_path)
        elif ext in {".png", ".jpg", ".jpeg"}:
            return self._parse_image(file_path)
        else:
            return [ParsedChunk(content=f"Unsupported file type: {ext}")]

    def _parse_image(self, file_path: str) -> List[ParsedChunk]:
        try:
            reader = _get_easyocr_reader()
            result = reader.readtext(file_path, detail=0)
            text = "\n".join(result)
            if text.strip():
                return [ParsedChunk(content=text, metadata={"source": "ocr"})]
        except ImportError:
            return [ParsedChunk(content=f"[Image OCR failed: easyocr not installed. Run: pip install easyocr]")]
        except Exception as e:
            return [ParsedChunk(content=f"[Image OCR error: {e}]")]
        return [ParsedChunk(content=f"[Image file: {Path(file_path).name} (OCR returned no text)]")]

    def _parse_pdf(self, file_path: str) -> List[ParsedChunk]:
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            chunks = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    chunks.append(ParsedChunk(content=text.strip(), page_number=i + 1))
            return chunks
        except Exception as e:
            return [ParsedChunk(content=f"PDF parse error: {e}")]

    def _parse_docx(self, file_path: str) -> List[ParsedChunk]:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return [ParsedChunk(content=text)] if text.strip() else [ParsedChunk(content="")]
        except Exception as e:
            return [ParsedChunk(content=f"DOCX parse error: {e}")]

    def _parse_pptx(self, file_path: str) -> List[ParsedChunk]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                if slide_texts:
                    texts.append("\n".join(slide_texts))
            return [ParsedChunk(content=t) for t in texts]
        except Exception as e:
            return [ParsedChunk(content=f"PPTX parse error: {e}")]

    def _parse_text(self, file_path: str) -> List[ParsedChunk]:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return [ParsedChunk(content=content)]

    def _parse_html(self, file_path: str) -> List[ParsedChunk]:
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            return [ParsedChunk(content=text)]
        except Exception as e:
            return [ParsedChunk(content=f"HTML parse error: {e}")]


def save_upload_file(file_bytes: bytes, filename: str) -> str:
    upload_dir = Path(settings.data_dir)
    upload_dir.mkdir(parents=True, exist_ok=True)

    stem = Path(filename).stem
    ext = Path(filename).suffix
    unique_name = f"{stem}_{uuid.uuid4().hex[:8]}{ext}"
    file_path = upload_dir / unique_name

    with open(file_path, "wb") as f:
        f.write(file_bytes)

    return str(file_path)


def compute_file_hash(file_bytes: bytes) -> str:
    return hashlib.sha256(file_bytes).hexdigest()
